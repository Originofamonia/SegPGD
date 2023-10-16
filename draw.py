"""
draw figures in paper
"""
import os
import re
from tqdm import tqdm
import pandas as pd
from sklearn.manifold import TSNE
from sklearn.decomposition import PCA
import seaborn as sns
import matplotlib.pyplot as plt
import torch
from pptx import Presentation
from pptx.util import Inches

# plt.rcParams.update({
#     "text.usetex": True,
#     "text.latex.preamble": r"\usepackage{amsfonts}"  # Load amsfonts package
# })

def tsne_hidden_layers():
    """
    1. use tsne to reduce dimension of each hidden layers
    2. plot hidden layers of clean vs adv
    """
    tsne = TSNE(n_components=2, perplexity=8, random_state=444)
    pca = PCA(n_components=10, random_state=444)
    labels = [1,2,3,4,5,6,7,8,9,10,11,12,13,14,15,16,1,2,3,4,5,6,7,8,9,10,11,12,13,14,15,16] # [1,2,1,2] or [1,2,3,4,1,2,3,4]
    layers = ['low_level', 'out']  # 'layer2', 'layer3',
    files = os.listdir(f'results/hidden_layers')
    batchs = [int(re.match(r'\d+', filename).group()) for filename in files]
    for i in tqdm(range(0, max(batchs), 4), desc=f'batches'):
        clean_features_1 = torch.load(f'results/hidden_layers/{i}_clean.pth')  # 1 means batch_1
        adv_features_1 = torch.load(f'results/hidden_layers/{i}_adv.pth')
        clean_features_2 = torch.load(f'results/hidden_layers/{i+1}_clean.pth')
        adv_features_2 = torch.load(f'results/hidden_layers/{i+1}_adv.pth')
        clean_features_3 = torch.load(f'results/hidden_layers/{i+2}_clean.pth')
        adv_features_3 = torch.load(f'results/hidden_layers/{i+2}_adv.pth')
        clean_features_4 = torch.load(f'results/hidden_layers/{i+3}_clean.pth')
        adv_features_4 = torch.load(f'results/hidden_layers/{i+3}_adv.pth')
        fig, axes = plt.subplots(1,len(layers))  # ,figsize=(14, 4)
        for j, l in enumerate(layers):
            clean_li = clean_features_1[l]
            adv_li = adv_features_1[l]
            clean_li_2 = clean_features_2[l]
            adv_li_2 = adv_features_2[l]
            clean_li_3 = clean_features_3[l]
            adv_li_3 = adv_features_3[l]
            clean_li_4 = clean_features_4[l]
            adv_li_4 = adv_features_4[l]
            cat_tensor = torch.cat(
                (clean_li,clean_li_2,clean_li_3,clean_li_4,adv_li,adv_li_2,adv_li_3,adv_li_4), dim=0).detach().cpu().numpy()
            cat_tensor = cat_tensor.reshape(cat_tensor.shape[0], -1)
            del clean_li, adv_li

            cat_tensor = pca.fit_transform(cat_tensor)
            x_tsne = tsne.fit_transform(cat_tensor)
            df = pd.DataFrame(x_tsne, columns=["comp1", "comp2"])
            df["y"] = labels

            sns_plot = sns.scatterplot(x="comp1", y="comp2", hue=df.y.tolist(),
                palette=sns.color_palette("hls", cat_tensor.shape[0]//2),
                data=df, ax=axes[j], legend=False)
            axes[j].set_xticks([])
            axes[j].set_yticks([])
            axes[j].set_xlabel('')
            axes[j].set_ylabel('')

        plt.savefig(f"results/visualize/t-SNE_{i}.png",bbox_inches='tight',pad_inches=0.13)
        plt.close(fig)


def compare_before_after_masks():
    """
    TODO: make a slides to compare before and after masks
    """
    before_files = os.listdir(f'results/before_AT')
    # after_files = os.listdir(f'results/after')
    pattern = r'^(\d+_\d+)'
    prefixes = set()

    for filename in before_files:
        match = re.match(pattern, filename)
        if match is not None:
            prefixes.add(match.group())

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]
    left = top = Inches(0.1)
    left_2 = Inches(6)
    table_top = Inches(6)
    width = Inches(14.0)
    height = Inches(1.2)
    prefixes = list(prefixes)

    for pre in tqdm(prefixes, desc=f'before_files: '):
        before = os.path.join('results/before_AT', f'{pre}_adv_y_pred.png')
        after = os.path.join('results/after', f'robust_{pre}_adv_y_pred.png')
        if os.path.exists(before) and os.path.exists(after):
            slide = prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(before, left, top)
            pic = slide.shapes.add_picture(after, left_2, top)

    prs.save(f'results/AT_before_after.pptx')


def prepare_5_columns():
    """
    draw: w/o defense, AT, AT+hloss
    """
    before_files = os.listdir(f'results/before_AT')
    # after_files = os.listdir(f'results/after')
    pattern = r'^(\d+_\d+)'
    prefixes = set()

    for filename in before_files:
        match = re.match(pattern, filename)
        if match is not None:
            prefixes.add(match.group())

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]
    left = top = Inches(0.1)
    left_2 = Inches(5)
    left_3 = Inches(10)
    top_2 = Inches(4)
    top_3 = Inches(8)
    width = Inches(4)
    height = Inches(2)
    prefixes = list(prefixes)

    for pre in tqdm(prefixes, desc=f'before_files: '):
        before = os.path.join('results/before_AT', f'{pre}_adv_y_pred.png')
        after = os.path.join('results/after', f'robust_{pre}_adv_y_pred.png')
        hloss = os.path.join('results/hloss', f'hloss_{pre}_adv_y_pred.png')
        hloss_clean = os.path.join('results/hloss', f'hloss_{pre}_clean_y_true.png')
        if os.path.exists(before) and os.path.exists(after):
            slide = prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(before, left, top)
            pic = slide.shapes.add_picture(after, left_2, top)
            pic = slide.shapes.add_picture(hloss, left_3, top)
            pic = slide.shapes.add_picture(hloss_clean, left, top_2)
            # Add a text box to the slide at the specified position
            text_box = slide.shapes.add_textbox(left, top_3, width, height)

            # Get the text frame within the text box
            tf = text_box.text_frame

            # Add a paragraph to the text frame
            p = tf.add_paragraph()
            p.text = pre

    prs.save(f'results/prepare_5_columns.pptx')


def selected_positive_negative():
    """
    selected good and negative examples:
    1. x_clean, 2. label, 3. w/o defense
    4. x_adv, 5. AT, 6. AT+hloss
    """
    positives = [
        '85_3','61_3','21_0','38-2','6_1','80_2','82_0','101_3',
        '115_3','16_3','58_1','68_3','75_3','53_1','72_3','91_3','0_0',
        '121_0','85_0']
    negatives = [
        '75_1','91_0','81_0','12_3','24_3','107_2']
    # relocation = {1: 'w/o defense', 2: 'AT', 3: 'AT+hloss', 4: 'clean'}
    reorder = {
        '85_3':'132','21_0':'132','6_1':'132','80_2':'312','82_0':'312',
        '101_3':'132','115_3':'132','16_3':'132','75_3':'132','72_3':'132',
        '91_3':'132','0_0':'132','121_0':'132','85_0':'132','91_0':'132',
        '81_0':'132','24_3':'132'}
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]
    left = top = Inches(0.1)
    left_2 = Inches(5)
    left_3 = Inches(10)
    top_2 = Inches(4)
    top_3 = Inches(8)
    width = Inches(4)
    height = Inches(2)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Positive examples"
    for pre in tqdm(positives, desc=f'before_files: '):
        before = os.path.join('results/before_AT', f'{pre}_adv_y_pred.png')
        after = os.path.join('results/after', f'robust_{pre}_adv_y_pred.png')
        hloss = os.path.join('results/hloss', f'hloss_{pre}_adv_y_pred.png')
        hloss_clean = os.path.join('results/hloss', f'hloss_{pre}_clean_y_true.png')
        four_images = [before, after, hloss, hloss_clean]  # select 3 our of 4
        if pre in reorder:
            three_images = [four_images[int(x) - 1] for x in reorder[pre]]
        else:
            three_images = four_images
        if os.path.exists(before) and os.path.exists(after):
            slide = prs.slides.add_slide(blank_slide_layout)
            fig, axs = plt.subplots(1, 3, figsize=(15, 6))
            before_img = plt.imread(three_images[0])  # w/o defense
            after_img = plt.imread(three_images[1])
            hloss_img = plt.imread(three_images[2])
            # hloss_clean_img = plt.imread(three_images[3])
            axs[0].imshow(before_img)
            axs[0].set_title(f'w/o defense')
            axs[0].axis('off')
            
            axs[1].imshow(after_img)
            axs[1].set_title(f'AT')
            axs[1].axis('off')

            axs[2].imshow(hloss_img)
            axs[2].set_title(f'AT+hloss')
            axs[2].axis('off')

            # axs[3].axis('off')
            # axs[3].imshow(hloss_clean_img)
            # axs[3].set_title(f'hloss_clean')

            img_filename = f'results/hloss/{pre}_overlay.png'
            fig.savefig(img_filename, bbox_inches='tight', pad_inches=0)
            plt.close()
            pic = slide.shapes.add_picture(img_filename, left, top)
            text_box = slide.shapes.add_textbox(left, top_3, width, height)

            # Get the text frame within the text box
            tf = text_box.text_frame

            # Add a paragraph to the text frame
            p = tf.add_paragraph()
            p.text = pre
    
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Negative examples"
    for pre in tqdm(negatives, desc=f'before_files: '):
        before = os.path.join('results/before_AT', f'{pre}_adv_y_pred.png')
        after = os.path.join('results/after', f'robust_{pre}_adv_y_pred.png')
        hloss = os.path.join('results/hloss', f'hloss_{pre}_adv_y_pred.png')
        hloss_clean = os.path.join('results/hloss', f'hloss_{pre}_clean_y_true.png')
        four_images = [before, after, hloss, hloss_clean]
        if pre in reorder:
            three_images = [four_images[int(x) - 1] for x in reorder[pre]]
        else:
            three_images = four_images
        if os.path.exists(before) and os.path.exists(after):
            slide = prs.slides.add_slide(blank_slide_layout)
            fig, axs = plt.subplots(1, 3, figsize=(15, 6))
            before_img = plt.imread(three_images[0])  # w/o defense
            after_img = plt.imread(three_images[1])
            hloss_img = plt.imread(three_images[2])
            # hloss_clean_img = plt.imread(three_images[3])
            axs[0].imshow(before_img)
            axs[0].set_title(f'w/o defense')
            axs[0].axis('off')
            
            axs[1].imshow(after_img)
            axs[1].set_title(f'AT')
            axs[1].axis('off')

            axs[2].imshow(hloss_img)
            axs[2].set_title(f'AT+hloss')
            axs[2].axis('off')

            # axs[3].imshow(hloss_clean_img)
            # axs[3].set_title(f'hloss_clean')
            # axs[3].axis('off')

            img_filename = f'results/hloss/{pre}_overlay.png'
            fig.savefig(img_filename, bbox_inches='tight', pad_inches=0)
            plt.close()
            pic = slide.shapes.add_picture(img_filename, left, top)
            text_box = slide.shapes.add_textbox(left, top_3, width, height)

            # Get the text frame within the text box
            tf = text_box.text_frame

            # Add a paragraph to the text frame
            p = tf.add_paragraph()
            p.text = pre

    prs.save(f'results/selected_5_columns.pptx')


def draw_qualitative_results():
    """
    3 positives: 80_2, 58_1, 92_3
    1 negative: 75_1
    clean image, adv image, perturbation, label
    w/o defense, after AT, after AT+hloss
    """
    selected_images = ['80_2', '58_1', '91_3', '75_1']
    fig, axs = plt.subplots(8, 4, figsize=(6.5, 10))  # figsize=(15, 6)
    reorder = {
        '85_3':'132','21_0':'132','6_1':'132','80_2':'312','82_0':'312',
        '101_3':'132','115_3':'132','16_3':'132','75_3':'132','72_3':'132',
        '91_3':'132','0_0':'132','121_0':'132','85_0':'132','91_0':'132',
        '81_0':'132','24_3':'132'}
    for i, pre in enumerate(selected_images):
        x_clean = plt.imread(f'results/qualitative_results/{pre}_x_clean.png')
        x_adv = plt.imread(f'results/qualitative_results/{pre}_x_adv.png')
        label = plt.imread(f'results/qualitative_results/{pre}_y_true.png')
        delta = plt.imread(f'results/qualitative_results/{pre}_delta.png')
        before = os.path.join('results/before_AT', f'{pre}_adv_y_pred.png')  # w/o defense
        after = os.path.join('results/after', f'robust_{pre}_adv_y_pred.png')  # AT
        hloss = os.path.join('results/hloss', f'hloss_{pre}_adv_y_pred.png')  # AT+hloss
        four_images = [before, after, hloss]  # hloss_clean
        if pre in reorder:
            three_images = [four_images[int(x) - 1] for x in reorder[pre]]
        else:
            three_images = four_images
        
        before_img = plt.imread(three_images[0])  # w/o defense
        after_img = plt.imread(three_images[1])
        hloss_img = plt.imread(three_images[2])
        axs[2*i, 0].imshow(x_clean)
        axs[2*i, 0].axis('off')
        axs[2*i, 1].imshow(x_adv)
        axs[2*i, 1].axis('off')
        axs[2*i, 2].imshow(label)
        axs[2*i, 2].axis('off')
        axs[2*i, 3].imshow(delta)
        axs[2*i, 3].axis('off')
        

        axs[2*i+1, 0].imshow(before_img)
        axs[2*i+1, 0].axis('off')
        axs[2*i+1, 1].imshow(after_img)
        axs[2*i+1, 1].axis('off')
        axs[2*i+1, 2].imshow(hloss_img)
        axs[2*i+1, 2].axis('off')
        axs[2*i+1, 3].axis('off')

        if i == 0:
            font = 9
            y = -0.2
            axs[2*i, 0].set_title(f'a', y=y, fontsize=font)
            axs[2*i, 1].set_title(f'b', y=y,fontsize=font)
            axs[2*i, 2].set_title(f'c', y=y,fontsize=font)
            axs[2*i, 3].set_title(f'd', y=y,fontsize=font)
            axs[2*i+1, 0].set_title(f'e', y=y,fontsize=font)
            axs[2*i+1, 1].set_title(f'f', y=y,fontsize=font)
            axs[2*i+1, 2].set_title(f'g', y=y,fontsize=font)
    
    img_filename = f'results/qualitative_results/qualitative.png'
    plt.tight_layout()
    fig.savefig(img_filename, bbox_inches='tight', pad_inches=0)
    plt.close()


def increasing_perturbations():
    epsilons = [0.005, 0.05, 0.1, 0.2, 0.3]
    fig, axs = plt.subplots(1, 5, figsize=(15, 5))
    for i, eps in enumerate(epsilons):
        img = plt.imread(f'results/clean/80_2_x_adv_{eps}.png')
        axs[i].imshow(img)
        axs[i].set_title(f'{eps}', y=-0.15,fontsize=19)
        axs[i].axis('off')
    
    img_filename = f'results/clean/epsilons.png'
    fig.suptitle('Epsilons', y=0.17,fontsize=19)
    plt.tight_layout()
    fig.savefig(img_filename, bbox_inches='tight', pad_inches=0)
    plt.close()


if __name__ == '__main__':
    # tsne_hidden_layers()
    # compare_before_after_masks()
    # prepare_5_columns()
    # selected_positive_negative()
    draw_qualitative_results()
    # increasing_perturbations()
